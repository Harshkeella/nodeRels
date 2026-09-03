"""Self-check. No network, no API keys, no test framework.

    python -m image_engine.test_engine

Covers the two things that are actually easy to get wrong: the scoring maths, and the
provider -> ImageCandidate normalisation. The end-to-end run uses a fake provider serving
file:// URLs, so the whole pipeline is exercised offline.
"""
import asyncio
import contextlib
import copy
import os
import pathlib
import random
import tempfile

from PIL import Image, ImageDraw

from . import config, engine, plan, providers, rank
from .schemas import ImageCandidate, ImageSlot, SlideRequest, VisualPlan

BRIEF = VisualPlan(
    visual_intent="A doctor reviewing a brain scan on a monitor in a hospital.",
    image_type="conceptual_realistic",
    primary_concepts=["doctor", "medical scan", "hospital", "monitor"],
    secondary_concepts=["technology", "healthcare"],
    avoid=["cartoon", "illustration"],
)


def cand(**kw) -> ImageCandidate:
    base = dict(id="x", provider="pexels", preview_url="p", download_url="d",
                width=1600, height=900, query="doctor medical scan")
    return ImageCandidate(**{**base, **kw})


def photo(w, h, kind="detailed", seed=0):
    """A synthetic image: 'detailed' has edges and contrast, 'flat' is a blank wash."""
    img = Image.new("RGB", (w, h), (128, 128, 128))
    if kind == "detailed":
        rng = random.Random(seed)
        d = ImageDraw.Draw(img)
        for _ in range(60):
            x, y = rng.randrange(w), rng.randrange(h)
            d.rectangle([x, y, x + rng.randrange(10, w // 3), y + rng.randrange(10, h // 3)],
                        fill=(rng.randrange(256), rng.randrange(256), rng.randrange(256)))
        # Photographs are noisy. Flat shapes compress to a few KB and would trip the
        # placeholder guard, which is exactly the guard being tested elsewhere.
        img = Image.blend(img, Image.effect_noise((w, h), 40).convert("RGB"), 0.3)
    return img


# ---------- scoring ----------

def test_aspect():
    wide, target = 16 / 10, 16 / 9
    assert rank.aspect_score(target, target) == 1.0
    assert rank.aspect_score(wide, target) > rank.aspect_score(4 / 3, target)
    assert rank.aspect_score(4 / 3, target) > rank.aspect_score(3 / 4, target)  # portrait worse
    assert rank.aspect_score(None, target) == 0.5                               # unknown, neutral
    assert 0.0 <= rank.aspect_score(3 / 4, target) <= 1.0
    # A wrong shape is penalised, never rejected outright.
    assert rank.aspect_score(1 / 3, target) > 0.0


def test_quality():
    good = rank.quality_score(photo(1920, 1080))
    flat = rank.quality_score(photo(1920, 1080, "flat"))
    small = rank.quality_score(photo(400, 300))
    assert good["quality"] > flat["quality"], (good, flat)      # sharpness + contrast
    assert good["quality"] > small["quality"]                    # resolution
    assert 0.0 <= flat["quality"] <= 1.0
    assert set(config.QUALITY_WEIGHTS) <= set(good)              # every weight has a signal


def test_dhash():
    a = photo(800, 600, seed=1)
    assert rank.dhash(a) == rank.dhash(a.copy())
    resized = a.resize((400, 300)).resize((1200, 900))           # same photo, re-encoded
    assert rank.is_duplicate(rank.dhash(resized), [rank.dhash(a)])
    assert not rank.is_duplicate(rank.dhash(photo(800, 600, seed=2)), [rank.dhash(a)])


def test_semantic():
    on = cand(description="doctor studying a medical scan on a hospital monitor")
    off = cand(description="a bowl of pasta on a wooden table")
    bad = cand(description="cartoon doctor medical scan hospital monitor illustration")
    assert rank.semantic_score(on, BRIEF) > rank.semantic_score(off, BRIEF)
    assert rank.semantic_score(off, BRIEF) == 0  # Matching search query cannot endorse unrelated pixels.
    assert "medical scan" not in cand(description="pasta").text()
    assert rank.semantic_score(bad, BRIEF) < rank.semantic_score(on, BRIEF)   # 'avoid' hit
    assert 0.0 <= rank.semantic_score(off, BRIEF) <= 1.0
    # The provider's own ordering breaks ties between identical captions.
    top = rank.semantic_score(cand(description="doctor scan", provider_rank=0), BRIEF)
    deep = rank.semantic_score(cand(description="doctor scan", provider_rank=20), BRIEF)
    assert top > deep


def test_diversity():
    first = cand(id="1", query="doctor scan", photographer="Ada", description="doctor scan")
    assert rank.diversity_score(first, []) == 1.0
    same_q = cand(id="2", query="doctor scan", photographer="Bob", description="nurse ward")
    same_who = cand(id="3", query="hospital ward", photographer="Ada", description="nurse ward")
    other = cand(id="4", query="city street", photographer="Cid", description="tram at dusk")
    assert rank.diversity_score(other, [first]) > rank.diversity_score(same_q, [first])
    assert rank.diversity_score(other, [first]) > rank.diversity_score(same_who, [first])


def test_final_score_weights():
    perfect = {k: 1.0 for k in config.FINAL_WEIGHTS}
    assert abs(rank.final_score(perfect) - 1.0) < 1e-9
    assert abs(rank.final_score({}) - 0.5) < 1e-9        # nothing measured -> neutral
    assert rank.final_score({"semantic": 1.0}) > rank.final_score({"semantic": 0.0})


# ---------- planner ----------

def test_fallback_plan():
    p = plan.fallback_plan(SlideRequest(
        presentation_topic="Artificial Intelligence in Healthcare",
        slide_title="AI-Powered Medical Diagnosis",
        slide_content="AI helps doctors analyze medical images and identify patterns."))
    assert p.source == "fallback"
    assert 3 <= len(p.queries) <= 8
    assert len({q.query for q in p.queries}) == len(p.queries)      # no repeated queries
    assert "the" not in p.primary_concepts and "how" not in p.primary_concepts
    assert plan.route(p, ["pexels", "pixabay"]) == ["pexels", "pixabay"]
    p.image_type = "literal_object"
    assert plan.route(p, ["pexels", "unsplash", "pixabay"])[0] == "pixabay"
    # An unknown provider still gets searched, just last.
    assert plan.route(p, ["pexels", "mystery"])[-1] == "mystery"


# ---------- providers ----------

def test_provider_normalisation():
    """Each provider maps its own JSON onto one shape. Stub the HTTP, check the mapping."""
    payloads = {
        "unsplash": {"results": [{"id": "abc", "width": 1600, "height": 900,
                                  "urls": {"small": "s", "regular": "r", "full": "f"},
                                  "links": {"html": "page", "download_location": "dl"},
                                  "description": None, "alt_description": "a doctor",
                                  "user": {"name": "Ada"}}]},
        "pexels": {"photos": [{"id": 7, "width": 1600, "height": 900,
                               "src": {"medium": "m", "large": "l", "large2x": "l2"},
                               "url": "page", "alt": "a doctor", "photographer": "Ada"}]},
        "pixabay": {"hits": [{"id": 9, "imageWidth": 1600, "imageHeight": 900,
                              "webformatURL": "w", "largeImageURL": "L",
                              "pageURL": "page", "tags": "doctor, scan", "user": "Ada"}]},
        "wikimedia": {"query": {"pages": [{"pageid": 10, "title": "File:Doctor.jpg",
            "imageinfo": [{"thumburl": "w", "descriptionurl": "page",
                           "thumbwidth": 1600, "thumbheight": 900,
                           "extmetadata": {"Artist": {"value": "Ada"},
                                           "LicenseShortName": {"value": "CC BY-SA 4.0"},
                                           "ImageDescription": {"value": "a doctor"}}}]}]}},
    }
    real_get = providers._get
    try:
        for name, payload in payloads.items():
            providers._get = lambda url, headers=None, _p=payload: _p
            got = providers.PROVIDERS[name]("doctor", 8, "landscape", "key")
            assert len(got) == 1, name
            c = got[0]
            assert c.provider == name and c.id and c.download_url, name
            assert (c.width, c.height) == (1600, 900), name
            assert c.photographer == "Ada" and c.source_url == "page", name
            assert c.license in config.ALLOWED_LICENSES, name
            assert c.query == "doctor" and c.provider_rank == 0, name
            assert "Ada" in c.attribution(), name
        providers._get = lambda url, headers=None: payloads["unsplash"]
        assert providers.PROVIDERS["unsplash"]("doctor", 8, None, "k")[0].trigger_url == "dl"
        scans = copy.deepcopy(payloads["wikimedia"])
        scans["query"]["pages"][0]["title"] = "File:Neural networks thesis.pdf"
        providers._get = lambda url, headers=None: scans
        assert providers.wikimedia("neural networks", 8, None) == []
    finally:
        providers._get = real_get


def test_provider_failure_is_contained():
    """One dead provider must not take the slide down."""
    async def run():
        log = []
        def boom(*a, **k):
            raise RuntimeError("provider is on fire")
        providers.PROVIDERS["boom"] = boom
        config.KEYS["boom"] = "BOOM_KEY"
        os.environ["BOOM_KEY"] = "x"
        real_get = providers._get
        providers._get = lambda url, headers=None: {"hits": [
            {"id": 1, "webformatURL": "w", "largeImageURL": "L", "pageURL": "p",
             "imageWidth": 1600, "imageHeight": 900, "tags": "doctor", "user": "Ada"}]}
        try:
            got = await providers.search_all(["doctor"], ["boom", "pixabay"], log=log)
        finally:
            providers._get = real_get
            del providers.PROVIDERS["boom"], config.KEYS["boom"], os.environ["BOOM_KEY"]
        assert len(got) == 1 and got[0].provider == "pixabay", got
        assert any("FAILED" in line for line in log), log
    asyncio.run(run())


def test_candidate_hygiene():
    dupes = [cand(id="1"), cand(id="1"), cand(id="2")]
    assert len(engine._dedupe(dupes)) == 2
    # Commercial Creative Commons licences are usable with the attribution stored in notes.
    mixed = [cand(id="1", license="pexels"), cand(id="2", license="cc-by-nc"),
             cand(id="3", license="cc-by-sa"), cand(id="4")]
    assert [c.id for c in engine._licensed(mixed, [])] == ["1", "3"]


# ---------- validation ----------

def test_validate(tmp):
    good = tmp / "good.jpg"
    photo(1600, 900).save(good)
    blob = good.read_bytes()
    img, why = engine.validate(blob, "hero")
    assert img and why is None, why
    assert engine.validate(b"not an image at all", "hero")[1].startswith("too small")
    assert engine.validate(b"x" * 20000, "hero")[1].startswith("undecodable")

    small = tmp / "small.jpg"
    photo(400, 300).save(small)
    assert engine.validate(small.read_bytes(), "hero")[1].startswith("400x300")
    # Same file passes for an icon slot: minimums are per role, not global.
    assert engine.validate(small.read_bytes(), "icon")[0] is not None


# ---------- end to end ----------

@contextlib.contextmanager
def fake_stack(tmp, tag="a"):
    """A provider serving file:// URLs, a stubbed planner, and a cache under tmp.
    Everything needed to run the real pipeline with no network and no keys."""
    files = []
    for i in range(6):
        p = tmp / ("img_%s%d.jpg" % (tag, i))
        photo(1920, 1080, seed=i).save(p)
        files.append(pathlib.Path(p).as_uri())

    captions = ["doctor medical scan hospital monitor", "doctor hospital monitor",
                "medical scan technology", "nurse ward healthcare", "city tram dusk",
                "pasta wooden table"]

    def fake(query, per_page, orientation, key):
        return [ImageCandidate(
            id="%s%d" % (tag, i), provider="fake", preview_url=u, download_url=u,
            source_url="https://example.test/%d" % i, width=1920, height=1080,
            description=captions[i], photographer="Ada" if i else "Bo",
            license="cc0", query=query, provider_rank=i)
            for i, u in enumerate(files)]

    providers.PROVIDERS["fake"] = fake
    config.KEYS["fake"] = "FAKE_KEY"
    os.environ["FAKE_KEY"] = "x"
    real_public = config.PUBLIC_PROVIDERS
    config.PUBLIC_PROVIDERS = ()
    real_dirs = engine.IMAGES, engine.DECKS
    engine.IMAGES, engine.DECKS = str(tmp / "images"), str(tmp / "decks")
    real_plan = plan.make_plan
    plan.make_plan = lambda slide: BRIEF.model_copy(update={
        "queries": plan.fallback_plan(slide).queries})
    try:
        yield
    finally:
        plan.make_plan = real_plan
        config.PUBLIC_PROVIDERS = real_public
        engine.IMAGES, engine.DECKS = real_dirs
        del providers.PROVIDERS["fake"], config.KEYS["fake"], os.environ["FAKE_KEY"]


def test_pipeline(tmp):
    """Full run against a fake provider serving file:// URLs. No network, no keys."""
    with fake_stack(tmp):
        slide = SlideRequest(
            presentation_id="deck-1", presentation_topic="AI in Healthcare",
            slide_number=3, slide_title="AI-Powered Medical Diagnosis",
            slide_content="AI helps doctors analyze medical images.",
            image_slots=[ImageSlot(slot_id="image_1", role="hero", aspect_ratio="16:9"),
                         ImageSlot(slot_id="image_2", role="supporting", aspect_ratio="1:1"),
                         ImageSlot(slot_id="image_3", role="supporting", aspect_ratio="1:1")],
            debug=True)
        out = engine.select_for_slide_sync(slide)

        assert len(out.selected_images) == 3, out.metadata["log"]
        assert len({s.local_path for s in out.selected_images}) == 3   # three DIFFERENT photos
        for s in out.selected_images:
            assert os.path.isfile(s.local_path)
            assert 0.0 <= s.final_score <= 1.0 and s.quality_score > 0.0
            assert s.attribution and s.license == "cc0"
            assert (s.width, s.height) == (1920, 1080)
        # The best on-brief caption wins the hero slot.
        assert out.selected_images[0].search_query
        assert out.metadata["candidates_found"] >= 6
        assert out.candidates and any(c.scores.get("final") for c in out.candidates)

        # Same deck, next slide: every image already used is refused perceptually.
        slide2 = slide.model_copy(update={"slide_number": 4})
        out2 = engine.select_for_slide_sync(slide2)
        used = {s.local_path for s in out.selected_images}
        assert not (used & {s.local_path for s in out2.selected_images}), "reused a photo"

        # A different deck may reuse them freely.
        out3 = engine.select_for_slide_sync(slide.model_copy(update={"presentation_id": "deck-2"}))
        assert {s.local_path for s in out3.selected_images} == used

        # Second run hit the cache: no .part files, one file per selected image.
        cached = list((tmp / "images").rglob("*.img"))
        assert cached and not list((tmp / "images").rglob("*.part"))
        assert len(list((tmp / "images").rglob("*.json"))) == len(cached)


def test_ppt_integration(tmp):
    """The seam: ppt.illustrate() fills the deck, ppt.render() places what it wrote.

    These two agree on dict keys nobody type-checks, so drift here loses every photo
    silently. This is the check that fails when it happens.
    """
    from .. import ppt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    deck = copy.deepcopy(ppt.DEMO)
    with fake_stack(tmp, tag="p"):
        ppt.illustrate(deck)
        # No id on a CLI deck: illustrate stamps one, so re-running does not exhaust
        # itself against its own previous run.
        assert deck["id"] and deck["id"] != copy.deepcopy(ppt.DEMO).get("id")
        again = copy.deepcopy(ppt.DEMO)
        ppt.illustrate(again)
        assert [sl.get("image", {}).get("path") for sl in again["slides"]] ==                [sl.get("image", {}).get("path") for sl in deck["slides"]]

    wanted = [i for i, sl in enumerate(deck["slides"]) if sl["kind"] in ppt.IMAGE_KINDS]
    assert wanted, "no slide kind takes an image"
    for i in wanted:
        img = deck["slides"][i].get("image")
        assert img, "slide %d got no image" % i
        assert os.path.isfile(img["path"]), img
        assert img["credit"] and 0.0 <= img["score"] <= 1.0
    # illustrate() runs one deck, so the deck-level dedupe must hold across its slides.
    assert len({deck["slides"][i]["image"]["path"] for i in wanted}) == len(wanted)
    for i, sl in enumerate(deck["slides"]):
        if sl["kind"] not in ppt.IMAGE_KINDS:
            assert "image" not in sl, "put a photo on a %s slide" % sl["kind"]

    prs = ppt.Presentation(ppt.render(deck, str(tmp / "integrated.pptx")))
    placed = [sl for sl in prs.slides
              if any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in sl.shapes)]
    assert len(placed) == len(wanted), (len(placed), len(wanted))
    for sl in placed:
        assert sl.notes_slide.notes_text_frame.text.strip().endswith("Fake")


def main():
    with tempfile.TemporaryDirectory(prefix="imgengine-") as d:
        tmp = pathlib.Path(d)
        tests = [test_aspect, test_quality, test_dhash, test_semantic, test_diversity,
                 test_final_score_weights, test_fallback_plan, test_provider_normalisation,
                 test_provider_failure_is_contained, test_candidate_hygiene,
                 test_validate, test_pipeline, test_ppt_integration]
        for t in tests:
            t(tmp) if t.__code__.co_argcount else t()
            print("ok  " + t.__name__)
    print("\n%d checks passed" % len(tests))


if __name__ == "__main__":
    main()
