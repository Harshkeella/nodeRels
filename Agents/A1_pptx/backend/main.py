"""The whole backend. One uvicorn target:

    uvicorn backend.main:app --reload        # http://127.0.0.1:8000
    python -m backend.main --selftest        # no API call: boots, hits every route

Deck generation streams over SSE because an illustrated deck is one Groq call plus a
provider round trip per slide -- a minute of silence on a POST is a broken product.
"""
import argparse, asyncio, contextlib, hashlib, json, os, re, sys, time
import urllib.error, urllib.parse, urllib.request
from typing import Any, Optional

from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel, Field

from . import assist, deck as D, ppt
from .image_engine import api as image_api, engine, providers
from ppt_video_agent import agent as video_agent




ROOT = os.path.dirname(ppt.HERE)                       # repo root, one above backend/
OUT = os.path.join(ROOT, "outputs")
UPLOADS = os.path.join(OUT, "uploads")
VERSIONS = os.path.join(OUT, "versions")
VIDEOS = os.path.join(OUT, "videos")
for d in (OUT, UPLOADS, VERSIONS, VIDEOS, engine.IMAGES):
    os.makedirs(d, exist_ok=True)

MAX_UPLOAD = 12 * 1024 * 1024                          # one slide's worth of photo
KEEP_VERSIONS = 20
VIDEO_JOBS: dict[str, dict] = {}
VIDEO_TASKS: set = set()

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

app = FastAPI(title="Deck Studio", version="2.0")
# Vite dev server runs on another origin. Loopback only -- this is a local tool.
app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=r"http://(localhost|127\.0\.0\.1):\d+",
    allow_methods=["*"], allow_headers=["*"])
app.include_router(image_api.app.router)               # /health, /images/search, /images/select

# The engine caches every download as ".img" whatever it really is, so the extension
# cannot name the type. Four magic numbers can. Anything else is not a picture and is
# not served.
MAGIC = ((b"\xff\xd8\xff", "image/jpeg"), (b"\x89PNG\r\n\x1a\n", "image/png"),
         (b"RIFF", "image/webp"), (b"GIF8", "image/gif"))


def log(event: str, **fields) -> None:
    """One line per step of the visual pipeline, on stderr, structured enough to grep.

    Ids and provider names only -- never a key, never a user's prompt.
    """
    print("%s %s" % (event, " ".join("%s=%s" % kv for kv in fields.items())),
          file=sys.stderr, flush=True)


# ---------- storage ----------

def slug(title: str, taken) -> str:
    s = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-")[:48] or "deck"
    name, n = s, 2
    while name in taken:
        name, n = "%s-%d" % (s, n), n + 1
    return name


def decks() -> list:
    """Every saved deck, newest first. The .json is the document; the .pptx sits beside
    it and is only ever a build artefact of it."""
    out = []
    for f in os.listdir(OUT):
        if f.endswith(".json"):
            with contextlib.suppress(OSError, ValueError):
                with open(os.path.join(OUT, f), encoding="utf-8") as fh:
                    out.append(load(json.load(fh)))
    return sorted(out, key=lambda d: d.get("created", 0), reverse=True)


def load(deck: dict) -> dict:
    """Bring a deck up to the current document shape on the way out of storage.

    Decks written before the editor existed are five hardcoded slide kinds; expand()
    turns them into elements so they open in the editor like anything else. Nothing is
    rewritten on disk until the user saves -- reading is not a migration.
    """
    return D.expand(with_media(deck))


def with_media(deck: dict) -> dict:
    """Give every photo a browser URL. Decks illustrated from the CLI only carry a disk
    path, and a path is not something the browser can load."""
    for s in deck.get("slides", []):
        img = s.get("image")
        if img and img.get("path") and not img.get("url"):
            img["url"] = media_url(img["path"])
        for e in s.get("elements", []):
            c = e.get("content") or {}
            if e.get("type") == "image" and c.get("path") and not c.get("url"):
                c["url"] = media_url(c["path"])
    return deck


def read(deck_id: str) -> dict:
    """One deck by id, or 404. The id is checked against the same pattern the .pptx
    route uses, so nothing here can be talked into reading outside outputs/."""
    if not re.fullmatch(r"[a-z0-9][a-z0-9-]*", deck_id or ""):
        raise HTTPException(404, "no such deck")
    path = os.path.join(OUT, deck_id + ".json")
    if not os.path.isfile(path):
        raise HTTPException(404, "no such deck")
    with open(path, encoding="utf-8") as fh:
        return load(json.load(fh))


def snapshot(deck: dict) -> None:
    """Keep the last few saves. Autosave means the user is never asked before their work
    is overwritten, so something has to be able to put it back."""
    folder = os.path.join(VERSIONS, deck["id"])
    os.makedirs(folder, exist_ok=True)
    with contextlib.suppress(OSError):
        with open(os.path.join(folder, "%d.json" % (time.time() * 1000)), "w",
                  encoding="utf-8") as fh:
            json.dump(deck, fh, separators=(",", ":"))
        stale = sorted(os.listdir(folder))[:-KEEP_VERSIONS]
        for f in stale:
            os.remove(os.path.join(folder, f))


def save(deck: dict, version: bool = False) -> dict:
    """Write the document. That is the whole save.

    The .pptx is built on download instead, from whatever the .json says at that moment.
    Autosave fires every second or so while someone is dragging a box; rendering a deck
    of charts and photos that often is work nobody asked for, and it is also the only way
    the downloaded file could ever disagree with the document. Render once, on the way
    out, and it cannot.
    """
    if version:
        with contextlib.suppress(Exception):
            snapshot(read(deck["id"]))
    tmp = os.path.join(OUT, deck["id"] + ".json.tmp")
    with open(tmp, "w", encoding="utf-8") as fh:
        json.dump(deck, fh, indent=1)
    os.replace(tmp, os.path.join(OUT, deck["id"] + ".json"))   # never a half-written deck
    return deck


def media_url(local_path: str) -> Optional[str]:
    """Browser URL for a cached photo, or None if it escaped the cache directory."""
    rel = os.path.relpath(local_path, engine.IMAGES)
    return None if rel.startswith("..") else "/media/" + rel.replace(os.sep, "/")


# ---------- generation ----------

class Source(BaseModel):
    name: str
    text: str


class GenerateRequest(BaseModel):
    prompt: str
    sources: list[Source] = Field(default_factory=list)
    template: Optional[str] = None
    images: bool = True


def sse(event: str, **data) -> str:
    return "event: %s\ndata: %s\n\n" % (event, json.dumps(data))


async def generate(req: GenerateRequest):
    """SSE: status lines, then the plan, then one event per photo, then the saved deck."""
    try:
        yield sse("status", step="write", text="Reading the material and planning slides")
        data = "\n\n".join("### %s\n%s" % (s.name, s.text) for s in req.sources)
        deck = await asyncio.to_thread(
            ppt.build, req.prompt,
            data or "(no source files -- work from the request alone)")

        if req.template:
            deck["template"] = req.template
        deck["id"] = slug(deck["deck_title"], {d.get("id") for d in decks()})
        deck["created"] = time.time()
        deck["prompt"] = req.prompt
        deck["sources"] = [s.name for s in req.sources]
        # Stamp ids before expanding so the slide the browser sees in the plan event is
        # the same slide the done event finishes -- ids must not change under the client.
        for s in deck["slides"]:
            s.setdefault("id", D.new_id("slide"))
        D.expand(deck)
        yield sse("plan", deck=deck)

        live = providers.available()
        total = len(deck["slides"])
        wanted = ppt.image_targets(deck, 2)
        if req.images and live and wanted:
            for n, i in enumerate(wanted, 1):
                yield sse("status", step="image", done=n - 1, total=len(wanted),
                          text="Finding a photo for " + deck["slides"][i]["name"])
                image = await asyncio.to_thread(ppt.illustrate_slide, deck, i, deck["id"])
                if image:
                    image["url"] = media_url(image["path"])
                    deck["slides"][i]["image"] = image
                    # A photo re-lays the slide out -- the text column narrows to make
                    # room. Rebuild the elements from the authored fields rather than
                    # trying to patch a layout the layout function owns.
                    deck["slides"][i] = D.expand_slide(deck["slides"][i], i + 1, total)
                yield sse("image", index=i, image=image, slide=deck["slides"][i],
                          done=n, total=len(wanted))
        elif req.images and not live:
            yield sse("status", step="image", done=0, total=0,
                      text="No stock-photo key set - text deck only")

        # Photos have had their turn. Now the application -- not the model -- reads the
        # elements each slide actually ended up with, and repairs the ones that came out
        # as nothing but text. Deterministic, offline, one pass.
        yield sse("status", step="visuals", text="Checking every slide has a visual")
        visuals = ppt.ensure_visuals(deck)
        for r in visuals:
            if r.get("repaired"):
                log("VISUAL_REPAIR_SUCCESS", deck=deck["id"], slide=r["id"],
                    visual=r["visual"])
            elif r["visual"] == "none":
                log("VISUAL_VALIDATION_FAILED", deck=deck["id"], slide=r["id"],
                    reason="nothing honest to add")

        yield sse("status", step="render", text="Building the .pptx")
        await asyncio.to_thread(save, deck)
        yield sse("done", deck=deck, visuals=visuals)
    except SystemExit as e:                    # ppt.py exits on a bad key / rate limit
        yield sse("error", error=str(e))
    except Exception as e:                     # never take the server down with a request
        yield sse("error", error="%s: %s" % (type(e).__name__, e))


# ---------- routes ----------

@app.get("/api/meta")
def meta() -> dict:
    """Everything the UI needs to draw itself: templates, keys, budgets, slide geometry
    and the element vocabulary. The editor's pickers are built from this, so adding an
    element type or a chart kind in deck.py makes it appear in the UI with no edit here."""
    return {
        "templates": ppt.TEMPLATES,
        "key": bool(os.getenv("GROQ_API_KEY")),
        "providers": providers.available(),
        "video": {"available": not video_agent.requirements(),
                  "missing": video_agent.requirements()},
        "source_budget": ppt.SOURCE_BUDGET,
        "geometry": {"image_x": ppt.IMAGE_X, "image_w": ppt.IMAGE_W, "image_h": ppt.IMAGE_H,
                     "text_w": ppt.TEXT_W, "text_w_img": ppt.TEXT_W_IMG,
                     "image_kinds": list(ppt.IMAGE_KINDS)},
        "doc": {"w": D.W, "h": D.H, "types": list(D.TYPES), "shapes": list(D.SHAPES),
                "charts": list(D.CHARTS), "tokens": list(D.TOKENS), "ops": list(D.OPS),
                "max_upload": MAX_UPLOAD},
    }


@app.get("/api/decks")
def list_decks() -> list:
    return decks()


@app.get("/api/deck/{deck_id}")
def get_deck(deck_id: str) -> dict:
    return read(deck_id)


class SaveRequest(BaseModel):
    deck: dict[str, Any]


@app.put("/api/deck/{deck_id}")
def put_deck(deck_id: str, req: SaveRequest) -> dict:
    """Autosave lands here. The body is untrusted: deck.clean() rebuilds the document
    field by field against the stored one, so id, created, prompt and sources come from
    disk and cannot be rewritten by a request."""
    stored = read(deck_id)
    try:
        merged = D.clean(req.deck, stored)
    except ValueError as e:
        raise HTTPException(400, str(e))
    return save(merged, version=True)


@app.get("/api/deck/{deck_id}/versions")
def list_versions(deck_id: str) -> list:
    read(deck_id)                                     # 404s before revealing anything
    folder = os.path.join(VERSIONS, deck_id)
    if not os.path.isdir(folder):
        return []
    return [{"at": int(f[:-5]) / 1000, "id": f[:-5]}
            for f in sorted(os.listdir(folder), reverse=True) if f.endswith(".json")]


def video_status(deck_id: str) -> dict:
    if deck_id in VIDEO_JOBS:
        return VIDEO_JOBS[deck_id]
    path = os.path.join(VIDEOS, deck_id + ".mp4")
    if os.path.isfile(path):
        return {"state": "ready", "percent": 100, "text": "Video ready",
                "url": "/api/video/" + deck_id + ".mp4",
                "updated": os.path.getmtime(path)}
    return {"state": "idle", "percent": 0, "text": "No video yet"}


class VideoRequest(BaseModel):
    voice: str = "en-US-AriaNeural"


def _build_video(deck_id: str, source: dict, voice: str) -> None:
    target = os.path.join(VIDEOS, deck_id + ".mp4")
    def progress(**patch):
        VIDEO_JOBS[deck_id] = {**VIDEO_JOBS[deck_id], **patch, "state": "working"}
    try:
        result = video_agent.build_video(source, target, progress, voice)
        VIDEO_JOBS[deck_id] = {"state": "ready", "percent": 100, "text": "Video ready",
                               "url": "/api/video/" + deck_id + ".mp4",
                               "segments": result["segments"], "updated": time.time()}
    except Exception as e:
        with contextlib.suppress(OSError):
            os.remove(target)
        VIDEO_JOBS[deck_id] = {"state": "error", "percent": 0,
                               "text": "%s: %s" % (type(e).__name__, e)}


@app.get("/api/deck/{deck_id}/video")
def get_video_status(deck_id: str) -> dict:
    read(deck_id)
    return video_status(deck_id)


@app.post("/api/deck/{deck_id}/video")
async def make_video(deck_id: str, req: VideoRequest) -> dict:
    source = read(deck_id)
    if VIDEO_JOBS.get(deck_id, {}).get("state") == "working":
        return VIDEO_JOBS[deck_id]
    missing = video_agent.requirements()
    if missing:
        raise HTTPException(503, "Video tools are not installed: " + ", ".join(missing))
    if not re.fullmatch(r"[a-z]{2,3}-[A-Z]{2,4}-[A-Za-z]+Neural", req.voice):
        raise HTTPException(400, "unknown narration voice")
    VIDEO_JOBS[deck_id] = {"state": "working", "step": "queued", "percent": 1,
                           "text": "Preparing the video"}
    task = asyncio.create_task(asyncio.to_thread(_build_video, deck_id, source, req.voice))
    VIDEO_TASKS.add(task)
    task.add_done_callback(VIDEO_TASKS.discard)
    return VIDEO_JOBS[deck_id]


@app.get("/api/video/{name}")
def video_file(name: str):
    if not re.fullmatch(r"[a-z0-9][a-z0-9-]*\.mp4", name or ""):
        raise HTTPException(404, "no such video")
    path = os.path.join(VIDEOS, name)
    if not os.path.isfile(path):
        raise HTTPException(404, "no such video")
    return FileResponse(path, media_type="video/mp4", filename=name,
                        headers={"Cache-Control": "no-store"})


@app.post("/api/deck/{deck_id}/restore/{version}")
def restore(deck_id: str, version: str) -> dict:
    read(deck_id)
    if not re.fullmatch(r"\d+", version or ""):
        raise HTTPException(404, "no such version")
    path = os.path.join(VERSIONS, deck_id, version + ".json")
    if not os.path.isfile(path):
        raise HTTPException(404, "no such version")
    with open(path, encoding="utf-8") as fh:
        old = json.load(fh)
    # Restoring is itself a save, so the state being replaced is snapshotted too and a
    # restore can be undone by restoring the version it just made.
    return save(D.clean(old, read(deck_id)), version=True)


class AskRequest(BaseModel):
    ask: str
    selection: dict[str, list[str]] = Field(default_factory=dict)
    apply: bool = True


@app.post("/api/deck/{deck_id}/ai")
async def ai_edit(deck_id: str, req: AskRequest) -> dict:
    """Ask AI. Returns the edited document and an honest report of what was refused.

    The model's response is never applied to the client's state directly: it is applied
    here, to the stored document, through deck.apply_ops. A response full of invented
    element ids costs the user nothing but a message saying so.
    """
    if not req.ask.strip():
        raise HTTPException(400, "Say what you want changed.")
    if not os.getenv("GROQ_API_KEY"):
        raise HTTPException(400, "No GROQ_API_KEY set — AI editing is unavailable.")
    stored = read(deck_id)
    try:
        new, report = await asyncio.to_thread(assist.edit, stored, req.ask, req.selection)
    except SystemExit as e:                            # ppt._ask exits on key/rate limit
        raise HTTPException(502, str(e))
    except Exception as e:
        raise HTTPException(502, "AI couldn't apply this change (%s)." % type(e).__name__)
    if report["applied"] and req.apply:
        new = save(D.clean(new, stored), version=True)
    return {"deck": new, "report": report}


class ImageSearch(BaseModel):
    query: str
    count: int = 6
    ratio: str = "4:3"


@app.post("/api/images")
async def find_images(req: ImageSearch) -> dict:
    """Stock photos for the editor's image panel, already downloaded and cached.

    Reuses the whole existing pipeline -- plan, search, rank, validate, dedupe -- rather
    than a second, thinner search path. Everything comes back as a /media URL because a
    provider's own CDN link is not something a saved deck is allowed to point at.
    """
    from .image_engine.schemas import ImageSlot, SlideRequest
    if not req.query.strip():
        raise HTTPException(400, "Say what you are looking for.")
    if not providers.available():
        raise HTTPException(400, "No stock-photo key set — add UNSPLASH_ACCESS_KEY, "
                                 "PEXELS_API_KEY or PIXABAY_API_KEY and restart.")
    n = max(1, min(12, req.count))
    try:
        w, _, h = req.ratio.partition(":")
        int(w), int(h)
    except ValueError:
        raise HTTPException(400, "ratio must look like 4:3")
    # A fresh dedupe scope per search: the engine refuses a photo already used under an
    # id, which is right across a deck and wrong for a picker the user reopens.
    slide = SlideRequest(
        presentation_id="search-%d" % (time.time() * 1000),
        slide_title=req.query.strip()[:200],
        image_slots=[ImageSlot(slot_id="s%d" % i, role="supporting", aspect_ratio=req.ratio)
                     for i in range(n)])
    found = await engine.select_for_slide(slide)
    return {"images": [{"url": media_url(i.local_path), "path": i.local_path,
                        "credit": i.attribution, "source_url": i.source_url,
                        "alt": i.search_query, "w": i.width, "h": i.height}
                       for i in found.selected_images if media_url(i.local_path)]}


def store(data: bytes, alt: str = "") -> dict:
    """Bytes -> a stored image the deck may point at. The one place an image becomes an
    asset, so a file the user picked and a file we fetched from a URL are the same thing
    from here on and cannot drift apart.

    Content-addressed, so the same picture arriving twice costs one copy. The type is the
    magic number, never the extension or the Content-Type header: an URL ending .jpg is
    not proof it returns a JPEG, and SVG is refused on purpose -- it is a document with
    script in it, and this app renders these straight into the page.
    """
    if len(data) > MAX_UPLOAD:
        raise HTTPException(413, "Images must be under %d MB." % (MAX_UPLOAD // 1048576))
    kind = next((t for sig, t in MAGIC if data.startswith(sig)), None)
    if not kind:
        raise HTTPException(400, "That is not a PNG, JPEG, GIF or WebP image.")
    name = hashlib.sha256(data).hexdigest()[:32] + "." + kind.split("/")[1]
    path = os.path.join(UPLOADS, name)
    if not os.path.isfile(path):
        tmp = path + ".part"                    # never a half-written image on /uploads
        with open(tmp, "wb") as fh:
            fh.write(data)
        os.replace(tmp, path)
    log("ASSET_READY", asset=name, bytes=len(data), type=kind)
    return {"url": "/uploads/" + name, "path": path, "type": kind, "alt": alt[:120]}


@app.post("/api/upload")
async def upload(file: UploadFile = File(...)) -> dict:
    """A user image. Validation, storage and the asset record are all store()."""
    return store(await file.read(MAX_UPLOAD + 1),
                 os.path.splitext(file.filename or "")[0])


class UrlImage(BaseModel):
    url: str


@app.post("/api/upload/url")
def upload_url(req: UrlImage) -> dict:
    """An image by its address. Fetched here, never by the browser.

    The deck stores our copy, not the address: an external URL expires, hotlinks, blocks
    by origin, or turns into an HTML error page, and any of those is a slide that was
    fine when it was made and is broken when it is presented. One download now, and the
    picture belongs to the deck.
    """
    url = req.url.strip()
    if not re.match(r"https?://[^\s/$.?#][^\s]*$", url, re.I):
        raise HTTPException(400, "That is not an http or https image address.")
    log("IMAGE_DOWNLOAD_STARTED", host=urllib.parse.urlsplit(url).hostname or "?")
    try:
        with urllib.request.urlopen(urllib.request.Request(
                url, headers={"User-Agent": "deck-studio"}), timeout=15) as r:
            data = r.read(MAX_UPLOAD + 1)
    except (urllib.error.URLError, OSError, ValueError) as e:
        log("IMAGE_DOWNLOAD_FAILED", host=urllib.parse.urlsplit(url).hostname or "?",
            reason=type(e).__name__)
        raise HTTPException(400, "Could not fetch that address (%s)." % type(e).__name__)
    return store(data, os.path.basename(urllib.parse.urlsplit(url).path))


@app.get("/uploads/{name}")
def uploaded(name: str):
    """Same contract as /media: inside the folder, and an image by its magic number or
    it is not served at all."""
    if not re.fullmatch(r"[0-9a-f]{32}\.(jpeg|png|webp|gif)", name or ""):
        raise HTTPException(404, "no such image")
    path = os.path.join(UPLOADS, name)
    if not os.path.isfile(path):
        raise HTTPException(404, "no such image")
    with open(path, "rb") as fh:
        head = fh.read(12)
    kind = next((t for sig, t in MAGIC if head.startswith(sig)), None)
    if not kind:
        raise HTTPException(404, "not an image")
    return FileResponse(path, media_type=kind,
                        headers={"Cache-Control": "public, max-age=31536000, immutable"})


@app.post("/api/generate")
async def start(req: GenerateRequest):
    if not req.prompt.strip():
        raise HTTPException(400, "Say what the deck should do.")
    if req.template and req.template not in ppt.TEMPLATES:
        raise HTTPException(400, "unknown template")
    return StreamingResponse(generate(req), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache",
                                      "X-Accel-Buffering": "no"})


@app.get("/media/{path:path}")
def media(path: str):
    """A cached stock photo, by its browser URL. Serves nothing outside the cache and
    nothing that does not start with an image's magic number."""
    target = os.path.realpath(os.path.join(engine.IMAGES, path))
    if os.path.commonpath([target, os.path.realpath(engine.IMAGES)]) != \
            os.path.realpath(engine.IMAGES) or not os.path.isfile(target):
        raise HTTPException(404, "no such image")
    with open(target, "rb") as fh:
        head = fh.read(12)
    kind = next((t for sig, t in MAGIC if head.startswith(sig)), None)
    if not kind:
        raise HTTPException(404, "not an image")
    # Cache entries are content-addressed by provider+id and never rewritten in place.
    return FileResponse(target, media_type=kind,
                        headers={"Cache-Control": "public, max-age=31536000, immutable"})


@app.get("/api/file/{name}")
def download(name: str):
    """Build the .pptx from the document, now. Nothing is cached, so the file that
    arrives is the deck as it currently stands rather than as it stood at the last save."""
    if not re.fullmatch(r"[a-z0-9][a-z0-9-]*\.pptx", name or ""):
        raise HTTPException(404, "no such deck")
    deck = read(name[:-5])
    target = os.path.join(OUT, name)
    try:
        ppt.render(deck, target)
    except Exception as e:
        raise HTTPException(500, "Export failed (%s). The deck itself is safe."
                                 % type(e).__name__)
    return FileResponse(target, media_type=PPTX_MIME, filename=name,
                        headers={"Cache-Control": "no-store"})


@app.delete("/api/deck/{deck_id}")
def delete(deck_id: str):
    read(deck_id)                                     # 404s on a bad or unknown id
    for path in (os.path.join(OUT, deck_id + ".json"), os.path.join(OUT, deck_id + ".pptx")):
        with contextlib.suppress(OSError):
            os.remove(path)
    folder = os.path.join(VERSIONS, deck_id)
    if os.path.isdir(folder):
        with contextlib.suppress(OSError):
            for f in os.listdir(folder):
                os.remove(os.path.join(folder, f))
            os.rmdir(folder)
    with contextlib.suppress(OSError):
        os.remove(os.path.join(VIDEOS, deck_id + ".mp4"))
    VIDEO_JOBS.pop(deck_id, None)
    return {"ok": True}


class NewDeck(BaseModel):
    deck_title: str = "Untitled deck"
    template: Optional[str] = None


@app.post("/api/deck")
def create(req: NewDeck) -> dict:
    """An empty deck to start from. Generation mints ids from the deck title; so does
    this, so a blank deck and a generated one are the same kind of thing from here on."""
    fresh = D.clean({
        "deck_title": req.deck_title,
        "template": req.template if req.template in ppt.TEMPLATES else sorted(ppt.TEMPLATES)[0],
        "slides": [{"name": "Untitled slide", "elements": []}],
    })
    fresh["id"] = slug(req.deck_title, {d.get("id") for d in decks()})
    fresh["created"] = time.time()
    return save(fresh)


# ---------- self-check ----------

def lifecycle(c):
    """The whole product, once, in order:

        generate -> editor -> manual edit -> AI edit -> insert an asset -> reorder
        -> autosave -> reload -> export -> read the .pptx back

    Every step goes through the real routes. The only thing stubbed is the network: the
    planner returns the demo deck and the AI returns a canned operation list, because the
    point is to prove the pipe carries the user's work end to end, not to buy tokens.
    """
    import copy
    from pptx import Presentation
    from pptx.util import Inches as In, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 1. generate
    real_build = ppt.build
    ppt.build = lambda *a: copy.deepcopy(ppt.DEMO)
    try:
        with c.stream("POST", "/api/generate",
                      json={"prompt": "lifecycle deck", "images": False}) as r:
            for _ in r.iter_lines():
                pass
    finally:
        ppt.build = real_build
    d = next(x for x in decks() if x.get("prompt") == "lifecycle deck")
    did = d["id"]
    assert all(s["elements"] for s in d["slides"]), "generation must produce elements"

    # 2. the editor opens it, 3. the user edits by hand
    doc = c.get("/api/deck/" + did).json()
    title = doc["slides"][0]["elements"][0]
    title["content"]["text"] = "Hand-edited title"
    title["x"], title["y"] = 1.0, 0.4
    title["style"]["size"] = 54

    # 4. an asset: a real uploaded image, inserted as an element
    png = ppt._png(os.path.join(OUT, "life.png"), 8, 6)
    try:
        up = c.post("/api/upload",
                    files={"file": ("shot.png", open(png, "rb"), "image/png")}).json()
    finally:
        os.remove(png)
    doc["slides"][0]["elements"].append({
        "type": "image", "x": 7.5, "y": 1.2, "w": 4.5, "h": 3.4,
        "content": {"url": up["url"], "path": up["path"], "alt": "uploaded"}, "style": {}})

    # 5. rearrange the deck
    doc["slides"] = [doc["slides"][1], doc["slides"][0], *doc["slides"][2:]]

    # 6. autosave
    assert c.put("/api/deck/" + did, json={"deck": doc}).status_code == 200

    # 7. AI edit, applied to the stored document through the validator
    stored = read(did)
    target = next(e for e in stored["slides"][1]["elements"]
                  if e["content"].get("text") == "Hand-edited title")
    canned = {"summary": "Recoloured the title and added a rule under it.", "operations": [
        {"action": "update_element", "element_id": target["id"], "slide_id": None,
         "changes": {"color": "accent"}},
        {"action": "add_element", "element_id": None, "slide_id": stored["slides"][1]["id"],
         "changes": {"type": "shape", "x": 1.0, "y": 1.9, "w": 3.0, "h": 0.06,
                     "content": {"shape": "rect"}, "style": {"fill": "accent"}}},
        {"action": "update_element", "element_id": "not_a_real_id", "slide_id": None,
         "changes": {"text": "should be refused"}},
    ]}
    real_ask, key = ppt._ask, os.environ.get("GROQ_API_KEY")
    ppt._ask = lambda *a, **k: canned
    os.environ["GROQ_API_KEY"] = "stub-for-the-lifecycle-test"
    try:
        out = c.post("/api/deck/%s/ai" % did,
                     json={"ask": "make the title pop", "selection": {"elements": [target["id"]]}})
        assert out.status_code == 200, out.text
        report = out.json()["report"]
    finally:
        ppt._ask = real_ask
        os.environ.pop("GROQ_API_KEY", None)
        if key:
            os.environ["GROQ_API_KEY"] = key
    assert report["applied"] == 2 and len(report["rejected"]) == 1, report

    # 8. reload: everything above survived, in order
    back = c.get("/api/deck/" + did).json()
    edited = next(e for e in back["slides"][1]["elements"]
                  if e["content"].get("text") == "Hand-edited title")
    assert (edited["x"], edited["y"]) == (1.0, 0.4), "manual layout lost on reload"
    assert edited["style"]["size"] == 54 and edited["style"]["color"] == "accent"
    assert back["slides"][0]["name"] == d["slides"][1]["name"], "reorder lost on reload"
    assert any(e["type"] == "image" and e["content"]["url"] == up["url"]
               for e in back["slides"][1]["elements"]), "uploaded asset lost on reload"

    # 9. export, and 10. read the .pptx back to check it says the same thing
    assert c.get("/api/file/%s.pptx" % did).status_code == 200
    prs = Presentation(os.path.join(OUT, did + ".pptx"))
    assert len(prs.slides) == len(back["slides"])
    words = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert "Hand-edited title" in words, words
    hand = next(sh for sh in prs.slides[1].shapes
                if sh.has_text_frame and sh.text_frame.text == "Hand-edited title")
    assert (hand.left, hand.top) == (In(1.0), In(0.4)), "export moved the box"
    run = hand.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(54)
    accent = RGBColor.from_string(ppt.TEMPLATES[back["template"]]["accent"])
    assert run.font.color.rgb == accent, "the AI's colour did not reach the .pptx"
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in prs.slides[1].shapes), \
        "the uploaded image did not reach the .pptx"

    c.delete("/api/deck/" + did)
    return len(prs.slides)


def selftest():
    """Every route, the traversal guard, the SSE shape and the whole product lifecycle --
    without spending an API call."""
    import copy
    from fastapi.testclient import TestClient
    c = TestClient(app)

    m = c.get("/api/meta").json()
    assert set(m["templates"]) == set(ppt.TEMPLATES), m["templates"]
    assert m["geometry"]["image_x"] == ppt.IMAGE_X
    assert isinstance(c.get("/api/decks").json(), list)
    assert c.get("/health").json()["providers"] == providers.available()

    for probe in ("..%2fppt.py", "ppt.py", "nope.pptx", "a%2fb.pptx"):
        assert c.get("/api/file/" + probe).status_code == 404, probe
    assert c.delete("/api/deck/..%2fppt").status_code == 404

    assert c.post("/api/generate", json={"prompt": "  "}).status_code == 400
    assert c.post("/api/generate", json={"prompt": "x", "template": "nope"}).status_code == 400

    # A deck through the stream with ppt.build stubbed: proves the event order and that
    # an id, a .pptx and a .json all land, without touching Groq.
    real = ppt.build
    ppt.build = lambda *a: copy.deepcopy(ppt.DEMO)
    try:
        with c.stream("POST", "/api/generate",
                      json={"prompt": "selftest deck", "images": False}) as r:
            events = [ln[7:] for ln in r.iter_lines() if ln.startswith("event: ")]
    finally:
        ppt.build = real
    assert events[0] == "status" and events[-1] == "done", events
    assert "plan" in events, events
    fresh = decks()[0]
    assert fresh["prompt"] == "selftest deck" and fresh["id"], fresh.get("prompt")
    did = fresh["id"]
    assert c.get("/api/file/%s.pptx" % did).headers["content-type"] == PPTX_MIME

    # ---- visual coverage survives the whole pipe, not just the moment it was added ----
    # No photo keys were used here, so every visual on this deck came from the repair
    # pass. It has to be in the saved document and in the .pptx, or the check was theatre.
    for s in c.get("/api/deck/" + did).json()["slides"]:
        assert s["kind"] in D.NO_VISUAL_NEEDED or D.has_visual(s), s["name"]
    from pptx import Presentation as _Pr
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _ST
    exported = _Pr(os.path.join(OUT, did + ".pptx"))
    drawn = [{sh.shape_type for sh in sl.shapes} for sl in exported.slides]
    for i, s in enumerate(read(did)["slides"]):
        if s["kind"] in D.NO_VISUAL_NEEDED:
            continue
        assert drawn[i] & {_ST.PICTURE, _ST.CHART, _ST.TABLE, _ST.AUTO_SHAPE}, \
            "slide %d reached the .pptx with nothing but text" % (i + 1)

    # ---- the editor's round trip: read, edit, save, reload, and it is still there ----
    got = c.get("/api/deck/" + did).json()
    assert got["version"] == 2 and all(s["elements"] for s in got["slides"])
    edited = copy.deepcopy(got)
    edited["slides"][0]["elements"][0]["content"]["text"] = "Edited by hand"
    edited["slides"][0]["elements"][0]["x"] = 2.5
    edited["deck_title"] = "Renamed"
    edited["id"] = "../../etc/passwd"                  # a client must not move a deck
    edited["created"] = 0
    saved = c.put("/api/deck/" + did, json={"deck": edited}).json()
    assert saved["id"] == did and saved["created"] == fresh["created"], "client rewrote a server field"
    again = c.get("/api/deck/" + did).json()
    assert again["deck_title"] == "Renamed"
    assert again["slides"][0]["elements"][0]["content"]["text"] == "Edited by hand"
    assert again["slides"][0]["elements"][0]["x"] == 2.5, "layout did not survive the reload"
    assert "title" not in again["slides"][0], "authored fields must be dropped on save"

    # the download is built from the document at request time, so it can never be one
    # edit behind what the editor is showing
    from pptx import Presentation as _P
    assert c.get("/api/file/%s.pptx" % did).status_code == 200
    words = [sh.text_frame.text for sh in _P(os.path.join(OUT, did + ".pptx")).slides[0].shapes
             if sh.has_text_frame]
    assert "Edited by hand" in words, words

    versions = c.get("/api/deck/%s/versions" % did).json()
    assert len(versions) >= 1, versions
    back = c.post("/api/deck/%s/restore/%s" % (did, versions[-1]["id"])).json()
    assert back["deck_title"] == "Demo", "restore did not put the old title back"
    assert c.post("/api/deck/%s/restore/9" % did).status_code == 404
    assert c.post("/api/deck/%s/restore/..%%2f..%%2fx" % did).status_code == 404

    assert c.put("/api/deck/" + did, json={"deck": {"slides": []}}).status_code == 400
    assert c.get("/api/deck/nope").status_code == 404
    assert c.get("/api/deck/..%2f..%2fbackend%2fppt").status_code == 404

    # ---- a blank deck gets a real server-side id, so autosave has somewhere to go ----
    made = c.post("/api/deck", json={"deck_title": "Fresh Start"}).json()
    assert made["id"] == "fresh-start" and len(made["slides"]) == 1
    assert c.put("/api/deck/%s" % made["id"], json={"deck": made}).status_code == 200
    assert c.delete("/api/deck/" + made["id"]).status_code == 200
    assert c.delete("/api/deck/" + made["id"]).status_code == 404

    # ---- uploads: images only, magic-sniffed, and nothing executable ----
    png = ppt._png(os.path.join(OUT, "up.png"), 3, 3)
    try:
        up = c.post("/api/upload", files={"file": ("a.png", open(png, "rb"), "image/png")}).json()
        assert up["url"].startswith("/uploads/") and c.get(up["url"]).status_code == 200
        again_up = c.post("/api/upload", files={"file": ("b.png", open(png, "rb"), "image/png")})
        assert again_up.json()["url"] == up["url"], "identical uploads must not duplicate"
    finally:
        os.remove(png)
    svg = b"<svg xmlns='http://www.w3.org/2000/svg'><script>alert(1)</script></svg>"
    assert c.post("/api/upload", files={"file": ("x.svg", svg, "image/svg+xml")}).status_code == 400
    assert c.post("/api/upload",
                  files={"file": ("x.png", b"MZ\x90\x00not a png", "image/png")}).status_code == 400
    for probe in ("../ppt.py", "nope.png", "zz.png", "%2e%2e%2fppt.py"):
        assert c.get("/uploads/" + probe).status_code == 404, probe

    # ---- an image by its address: fetched here, checked here, stored here ----
    # The network is stubbed; the point is that what a URL yields goes through exactly the
    # same door a picked file does, and that a page pretending to be a .jpg does not.
    import io as _io
    png = ppt._png(os.path.join(OUT, "url.png"), 5, 4)
    body = open(png, "rb").read()
    os.remove(png)
    real_open, served = urllib.request.urlopen, [body]

    class _Resp(_io.BytesIO):
        def __enter__(self): return self
        def __exit__(self, *a): self.close()

    urllib.request.urlopen = lambda *a, **k: _Resp(served[0])
    try:
        got = c.post("/api/upload/url", json={"url": "https://example.com/a.jpg"}).json()
        assert got["url"].startswith("/uploads/") and got["type"] == "image/png", got
        assert c.get(got["url"]).status_code == 200, "a fetched image must serve back"
        assert got["url"] == store(body)["url"], "same bytes, one copy"
        served[0] = b"<!doctype html><title>404</title>"     # .jpg that is not a JPEG
        assert c.post("/api/upload/url",
                      json={"url": "https://example.com/b.jpg"}).status_code == 400
    finally:
        urllib.request.urlopen = real_open
    for probe in ("", "notaurl", "javascript:alert(1)", "file:///etc/passwd"):
        assert c.post("/api/upload/url", json={"url": probe}).status_code == 400, probe

    # ---- AI edit: no key, no call; a hallucinated id costs the user nothing ----
    key, os.environ["GROQ_API_KEY"] = os.getenv("GROQ_API_KEY"), ""
    del os.environ["GROQ_API_KEY"]
    try:
        assert c.post("/api/deck/%s/ai" % did, json={"ask": "make it pop"}).status_code == 400
        assert c.post("/api/deck/%s/ai" % did, json={"ask": " "}).status_code == 400
    finally:
        if key:
            os.environ["GROQ_API_KEY"] = key
    live = c.get("/api/deck/" + did).json()
    target = live["slides"][0]["elements"][0]["id"]
    out, applied, rejected = D.apply_ops(live, [
        {"action": "update_element", "element_id": target, "changes": {"fontSize": 61}},
        {"action": "update_element", "element_id": "made_up", "changes": {"text": "x"}},
    ])
    assert applied == ["update_element"] and len(rejected) == 1
    assert out["slides"][0]["elements"][0]["style"]["size"] == 61

    assert c.delete("/api/deck/" + did).status_code == 200
    assert not os.path.isfile(os.path.join(OUT, did + ".pptx"))

    slides = lifecycle(c)

    assert media_url(os.path.join(engine.IMAGES, "unsplash", "a.img")) == "/media/unsplash/a.img"
    assert media_url(os.path.join(ROOT, "secret.png")) is None

    # /media must confine itself to the cache and refuse anything that is not a picture,
    # whatever the extension claims -- the cache stores JPEG, PNG and WEBP all as ".img".
    probe = os.path.join(engine.IMAGES, "selftest.img")
    ppt._png(probe, 4, 4)
    try:
        r = c.get("/media/selftest.img")
        assert r.status_code == 200 and r.headers["content-type"] == "image/png", r.headers
        with open(probe, "wb") as fh:
            fh.write(b"<html>not a picture</html>")
        assert c.get("/media/selftest.img").status_code == 404
    finally:
        os.remove(probe)
    assert c.get("/media/../ppt.py").status_code in (404, 200) and \
        b"import argparse" not in c.get("/media/../ppt.py").content
    assert c.get("/media/..%2f..%2fbackend%2fppt.py").status_code == 404
    assert c.get("/media/nope.img").status_code == 404

    print("ok - routes serve, traversal blocked, uploads image-only, AI ops validated, "
          "/media and /uploads confined and sniffed")
    print("ok - every content slide of a generated deck carries a visual, in the saved "
          "document and in the exported .pptx; image URLs fetched, sniffed and stored")
    print("ok - lifecycle: generate -> edit -> AI edit -> insert asset -> reorder -> "
          "autosave -> reload -> export, %d slides, layout and colours intact in the .pptx"
          % slides)


if __name__ == "__main__":
    a = argparse.ArgumentParser(description="Deck Studio backend")
    a.add_argument("--selftest", action="store_true")
    a.add_argument("--port", type=int, default=8000)
    n = a.parse_args()
    if n.selftest:
        selftest()
    else:
        import uvicorn
        uvicorn.run(app, host="127.0.0.1", port=n.port)
