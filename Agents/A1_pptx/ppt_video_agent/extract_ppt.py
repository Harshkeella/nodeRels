from pptx import Presentation
import json
import os


import glob

ppt_files = glob.glob("input/*.pptx")

if not ppt_files:
    raise FileNotFoundError("No .pptx file found inside the input folder.")

PPT_PATH = ppt_files[0]
OUTPUT_PATH = "output/presentation_data.json"


def extract_presentation(ppt_path):

    presentation = Presentation(ppt_path)

    presentation_data = []

    for slide_number, slide in enumerate(presentation.slides, start=1):

        slide_data = {
            "slide_number": slide_number,
            "title": "",
            "text": [],
            "tables": []
        }

        # -------------------------
        # Extract slide title
        # -------------------------

        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()

        # -------------------------
        # Extract text
        # -------------------------

        for shape in slide.shapes:

            if shape.has_text_frame:

                text = shape.text.strip()

                if text:
                    # Avoid repeating title
                    if text != slide_data["title"]:
                        slide_data["text"].append(text)

            # -------------------------
            # Extract tables
            # -------------------------

            if shape.has_table:

                table_data = []

                for row in shape.table.rows:

                    row_data = []

                    for cell in row.cells:
                        row_data.append(cell.text.strip())

                    table_data.append(row_data)

                slide_data["tables"].append(table_data)

        presentation_data.append(slide_data)

    return presentation_data


def save_json(data, output_path):

    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with open(output_path, "w", encoding="utf-8") as file:
        json.dump(
            data,
            file,
            indent=4,
            ensure_ascii=False
        )


if __name__ == "__main__":

    print("Reading PowerPoint...")

    data = extract_presentation(PPT_PATH)

    save_json(data, OUTPUT_PATH)

    print(f"Slides extracted: {len(data)}")
    print(f"Data saved to: {OUTPUT_PATH}")

    print("\n--- Extracted Content ---")

    for slide in data:

        print(f"\nSlide {slide['slide_number']}")
        print(f"Title: {slide['title']}")

        for text in slide["text"]:
            print(f"Text: {text}")

        for table in slide["tables"]:
            print(f"Table: {table}")