"""Offline regression check: real MCP transport, rendering, tenant boundaries, and editor round trip.

Run from the repo root after installing both services and packages/artifact-core:
    python scripts/verify_artifacts.py
No model calls, API keys, or existing user files are used.
"""
import asyncio
import copy
import json
import os
from pathlib import Path
import socket
import subprocess
import sys
import tempfile
import time
import uuid
from unittest.mock import patch, AsyncMock

ROOT = Path(__file__).resolve().parents[1]
sys.path[:0] = [str(ROOT / "backend"), str(ROOT / "Agents/A1_pptx")]
os.environ["AGENT_SHARED_SECRET"] = "artifact-regression-secret-not-for-deployment-12345"
os.environ["GROQ_API_KEY"] = ""
os.environ["OPENROUTER_API_KEY"] = ""
os.environ["OPIK_TRACK_DISABLE"] = "true"
os.environ["PYTHON_DOTENV_DISABLED"] = "1"

import fitz
import httpx
from pptx import Presentation
from noderels_artifacts import Block, Document, Section, from_markdown
from noderels_artifacts.pdf import render_pdf
from noderels_artifacts.jobs import Jobs
from backend import ppt, grounded
from ppt_video_agent import agent as video
from app.core import auth
from app.services import agent_client, artifacts, chat_store

SAMPLE = """# Retrieval architecture

## How retrieval works
Hybrid retrieval combines vector similarity with graph relationships.
The reranker selects the most relevant source chunks.

## Measured data
| Model | Accuracy |
| --- | --- |
| Baseline | 82.75% |
| Hybrid | 94.125% |

## Formula
$$
E = mc^2
$$
"""


async def verify(root: Path):
    document = from_markdown(SAMPLE, "Brief", ["source.md"])
    assert document.title == "Retrieval architecture"
    assert document.sections[1].blocks[0].rows[2][1] == "94.125%"
    assert document.sections[2].blocks[0].text == "E = mc^2"
    assert from_markdown("## Topic\nText.\n### Detail\nMore.", "Brief", []).sections[1].level == 3
    # The renderer numbers pages itself, so a model-written "Slide 3:" prefix is stripped.
    numbered = from_markdown("## Slide 3: Results\nText.\n## Section 1 - Aside\nMore.", "Brief", [])
    assert [s.title for s in numbered.sections] == ["Results", "Aside"]
    assert from_markdown("## Q3 2024 revenue\nText.", "Brief", []).sections[0].title == "Q3 2024 revenue"
    # The title names the subject, never the deliverable the model was asked for.
    assert from_markdown("# Video Script: ETL\nText.", "B", []).title == "ETL"
    assert from_markdown("# Presentation - Q3 Results\nText.", "B", []).title == "Q3 Results"
    assert from_markdown("# ETL: A Practical Guide\nText.", "B", []).title == "ETL: A Practical Guide"
    assert from_markdown("# Documentation Standards\nText.", "B", []).title == "Documentation Standards"
    assert from_markdown("# Topic\n![visual](https://example.com/a.png)\n```\n```\n## Body\nSource words.", "B", []).sections[0].blocks[0].text == "Source words."
    pdf = root / "sample.pdf"
    render_pdf(document, pdf)
    with fitz.open(pdf) as rendered:
        assert "94.125%" in "".join(p.get_text() for p in rendered)
        rendered[0].get_pixmap(matrix=fitz.Matrix(1.4, 1.4)).save(root / "sample.png")
        for page in rendered:
            for x0, y0, x1, y1, *_ in page.get_text("blocks"):
                assert x0 >= 45 and x1 <= page.rect.width - 45 and y1 < page.rect.height - 15

    queue = Jobs(root / "queue-check")
    job = queue.submit("alice", {"format": "pdf"})
    assert queue.submit("alice", {"format": "pdf"}, job_id=job["id"])["id"] == job["id"]
    try:
        queue.get("bob", job["id"])
        raise AssertionError("cross-user read succeeded")
    except KeyError:
        pass
    assert queue.claim()["id"] == job["id"] and queue.claim() is None
    assert Jobs(root / "queue-check").get("alice", job["id"])["state"] == "running"
    queue.update(job["id"], state="done")

    for question, expected in [("Can you generate a PDF about retrieval?", "pdf"),
                               ("Create slides about retrieval", "pptx"),
                               ("Make a video of this PPT", "video"),
                               ("How can I generate a PDF?", None),
                               ("Do not create a video", None)]:
        assert artifacts.intent(question) == expected, question
    assert artifacts.intent("I want you to generate videos about neural networks") == "video"
    assert artifacts.intent("Please create presentations") == "pptx"
    assert artifacts.intent("Generate PDFs") == "pdf"
    assert artifacts.intent("Create a PDF about what is hybrid retrieval") == "pdf"

    points = [{"slide_number": 1, "text": "Accuracy is 94.125%.", "notes": ""}]
    original = ppt._ask
    ppt._ask = lambda *a, **k: (_ for _ in ()).throw(SystemExit("offline"))
    try:
        assert video.narrate({}, points) == ["Accuracy is 94.125%."]
        for bad in ["Let's focus on this point. Accuracy is 94.125%.", "Accuracy is 98.1%."]:
            ppt._ask = lambda *a, **k: {"slides": [{"slide_number": 1, "narrations": [bad]}]}
            assert video.narrate({}, points) == ["Accuracy is 94.125%."]
    finally:
        ppt._ask = original

    # The Studio planner controls design; photos are copied into the job, not linked to its cache.
    fixture = root / "photo.png"
    ppt._png(str(fixture), 800, 1000)
    designed = root / "designed"
    designed.mkdir()
    with patch.object(ppt, "_ask", return_value={"template": "academic", "photo_sections": [0]}) as planner, \
         patch.object(ppt, "illustrate_slide", return_value={"path": str(fixture), "query": "retrieval", "credit": "Fixture"}):
        styled = grounded.build(document, designed, str(uuid.uuid4()), "Use academic styling")
    assert styled["template"] == "academic" and planner.call_args.kwargs["name"] == "grounded_deck"
    assert "Use academic styling" in planner.call_args.args[0][1]["content"]
    assert all(Path(e["content"]["path"]).parent == designed for s in styled["slides"] for e in s["elements"] if e["type"] == "image")
    ppt.render(styled, str(designed / "deck.pptx"), strict=True)
    wrapped = {"slides": [{"elements": [{"type": "text", "w": 8, "h": 1, "content": {"text": "One sentence\nwrapped for layout."}},
                                        {"type": "image", "content": {"path": "image-cover.png", "alt": "stock search"}}]}]}
    assert [p["text"] for p in video.speaking_points(wrapped)] == ["One sentence wrapped for layout."]
    assert "w=1:h=1" in video._highlight({"x": 50, "y": 50, "w": 1, "h": 1}, {})
    matrix = Document(title="Formula", sections=[Section(title="Matrix", blocks=[Block(kind="equation", text=r"\begin{pmatrix}1 & 2\\3 & 4\end{pmatrix}")])])
    render_pdf(matrix, root / "matrix.pdf")
    with fitz.open(root / "matrix.pdf") as rendered:
        assert "pmatrix" in rendered[0].get_text()
    with patch.object(ppt, "_ask", side_effect=SystemExit("offline")), patch.object(ppt, "illustrate_slide", return_value=None):
        matrix_deck = grounded.build(matrix, root, str(uuid.uuid4()))
    assert any("pmatrix" in e["content"].get("text", "") for s in matrix_deck["slides"] for e in s["elements"])

    # Transient TTS failure retries. Repeated failure records all lines locally.
    import types
    speech_folder = root / "speech"
    speech_folder.mkdir()
    saves = 0
    async def save(path):
        nonlocal saves
        saves += 1
        if saves == 1:
            raise RuntimeError("transient speech failure")
        Path(path).write_bytes(b"audio")
    edge = types.SimpleNamespace(Communicate=lambda *a: types.SimpleNamespace(save=save))
    with patch.dict(sys.modules, {"edge_tts": edge}), patch.object(asyncio, "sleep", new=AsyncMock()):
        assert len(await video._speak_all(["Hello"], speech_folder, "en-US-AriaNeural")) == 1 and saves == 2
        edge.Communicate = lambda *a: types.SimpleNamespace(save=AsyncMock(side_effect=RuntimeError("offline")))
        with patch.object(video, "_speak_local", return_value=[speech_folder / "local.wav"]) as local:
            assert (await video._speak_all(["Hello"], speech_folder, "en-US-AriaNeural"))[0].name == "local.wav"
            local.assert_called_once()

    with socket.socket() as sock:
        sock.bind(("127.0.0.1", 0))
        port = sock.getsockname()[1]
    url = f"http://127.0.0.1:{port}"
    agent_client.registry.cache_clear()
    os.environ["DECK_AGENT_URL"] = url
    environment = {**os.environ, "AGENT_STORAGE_DIR": str(root / "agent"), "IMG_CACHE": str(root / "images"), "PYTHONUNBUFFERED": "1"}
    # The agent runs on its own venv when it has one, so this check exercises the
    # environment the agent actually ships with rather than the runner's.
    agent_python = next((str(p) for p in (ROOT / "Agents/A1_pptx/.venv/Scripts/python.exe",
                                          ROOT / "Agents/A1_pptx/.venv/bin/python") if p.exists()), sys.executable)
    # Disable providers in this test process only; the offline check makes no network searches.
    bootstrap = ("from backend.image_engine import providers; providers.available = lambda: []; "
                 f"import uvicorn; uvicorn.run('backend.integrated:app', host='127.0.0.1', port={port}, access_log=False)")
    process = subprocess.Popen([agent_python, "-c", bootstrap],
                               cwd=ROOT / "Agents/A1_pptx", env=environment,
                               stdout=(root / "agent.log").open("w"), stderr=subprocess.STDOUT,
                               creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
    try:
        async with httpx.AsyncClient(timeout=2) as client:
            for _ in range(60):
                try:
                    if (await client.get(url + "/health")).status_code == 200:
                        break
                except httpx.HTTPError:
                    pass
                await asyncio.sleep(.5)
            else:
                raise RuntimeError((root / "agent.log").read_text())
            assert (await client.post(url + "/mcp", json={})).status_code == 401

        auth.run_as(auth.User("alice"))
        job_id = str(uuid.uuid4())
        submitted = await agent_client.call("submit_artifact", {"job_id": job_id, "document": document.model_dump()})
        assert submitted["id"] == job_id, submitted
        for _ in range(60):
            status = await agent_client.call("artifact_status", {"job_id": job_id})
            if status["state"] in ("done", "failed"):
                break
            await asyncio.sleep(.5)
        assert status["state"] == "done", status
        deck = await agent_client.call("studio_document", {"job_id": job_id})
        assert len(deck["slides"]) == 4
        assert deck["slides"][2]["elements"][1]["content"]["rows"][2][1] == "94.125%"
        words = [shape.text for slide in Presentation(root / "agent" / job_id / "deck.pptx").slides for shape in slide.shapes if shape.has_text_frame]
        assert any("Hybrid retrieval combines" in word for word in words)
        edited = copy.deepcopy(deck)
        edited["deck_title"] = "Edited title"
        edited["slides"][1]["elements"][1]["content"]["text"] = "A user-authored edit."
        edited["slides"][2]["elements"][1]["content"]["rows"][1][0] = "Detailed table cell " * 15
        duplicate = copy.deepcopy(edited["slides"][3]["elements"][1])
        duplicate["id"] = "copied-equation"
        edited["slides"][3]["elements"].append(duplicate)
        saved = await agent_client.call("save_studio_document", {"job_id": job_id, "document": edited})
        assert saved["revision"] == 1 and saved["deck_title"] == "Edited title"
        assert saved["slides"][2]["elements"][1]["content"]["rows"][1][0] == "Detailed table cell " * 15
        assert sum(e["type"] == "image" for e in saved["slides"][3]["elements"]) == 2
        try:
            await agent_client.call("save_studio_document", {"job_id": job_id, "document": edited})
            raise AssertionError("stale revision was accepted")
        except ValueError as exc:
            assert "changed in another tab" in str(exc)
        auth.run_as(auth.User("bob"))
        try:
            await agent_client.call("studio_document", {"job_id": job_id})
            raise AssertionError("cross-user MCP read succeeded")
        except ValueError as exc:
            assert "Artifact not found" in str(exc)
        async with httpx.AsyncClient() as client:
            assert (await client.get(url + f"/files/{job_id}/deck.pptx", headers=agent_client.headers())).status_code == 404
        auth.run_as(auth.User("alice"))
        async with httpx.AsyncClient() as client:
            response = await client.get(url + f"/files/{job_id}/deck.pptx", headers={**agent_client.headers(), "Range": "bytes=0-99"})
            assert response.status_code == 206 and len(response.content) == 100
            photo = await client.get(url + f"/files/{job_id}/image-cover.png", headers=agent_client.headers())
            assert photo.status_code == 200 and photo.content.startswith(b"\x89PNG")
            assert (await client.get(url + f"/files/{job_id}/request.json", headers=agent_client.headers())).status_code == 404
        if "--video" in sys.argv:
            # Real PowerPoint/LibreOffice + local speech + FFmpeg, no speech/model network dependency.
            target = root / "video's output.mp4"
            source = root / "video-deck.json"
            source.write_text(json.dumps({**styled, "slides": [styled["slides"][1]]}), encoding="utf-8")
            code = """import json, subprocess, sys
from pathlib import Path
from ppt_video_agent import agent
agent.build_video(json.loads(Path(sys.argv[1]).read_text(encoding='utf-8')), sys.argv[2])
probe = subprocess.run([agent._ffmpeg(), '-i', sys.argv[2], '-f', 'null', '-'], capture_output=True, text=True)
assert probe.returncode == 0 and 'Audio:' in probe.stderr and 'Video:' in probe.stderr, probe.stderr
print('PASS: real slide rendering, local narration and decoded MP4 audio/video')
"""
            result = await asyncio.to_thread(subprocess.run, [agent_python, "-c", code, str(source), str(target)],
                cwd=ROOT / "Agents/A1_pptx", env={**environment, "VIDEO_TTS_ENGINE": "local"},
                capture_output=True, text=True, timeout=600,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0)
            assert result.returncode == 0, result.stdout + result.stderr
            assert target.stat().st_size > 1000
            print(result.stdout.strip())
        print("PASS: PDF layout/text/formula fallback, Studio planner/images, exact tables, durable queue, MCP transport, tenant isolation, file ranges, editor revisions, speech retries/fallback")
    finally:
        process.terminate()
        process.wait(timeout=15)


if __name__ == "__main__":
    destination = Path(sys.argv[1]).resolve() if len(sys.argv) > 1 else Path(tempfile.mkdtemp(prefix="noderels-artifacts-"))
    destination.mkdir(parents=True, exist_ok=True)
    asyncio.run(verify(destination))
    print("Verification files:", destination)
